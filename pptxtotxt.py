from pptx import Presentation

def extract_text_from_pptx(file_path):
    """
    从 .pptx 文件中提取所有文本内容。

    参数:
        file_path (str): PPTX 文件的路径。

    返回:
        str: 所有幻灯片中文本内容拼接后的字符串，每段文本用换行符分隔。
    """
    try:
        prs = Presentation(file_path)
        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        
        return "\n".join(text_runs)
    
    except Exception as e:
        raise RuntimeError(f"无法读取PPTX文件 {file_path}: {e}")
    

# text = extract_text_from_pptx(r"E:\桌面\whut\深入学习贯彻党的二十届三中全会精神.pptx")
# print(text)